from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
NAVY = RGBColor(0x1A, 0x23, 0x7E)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)

def add_slide_number(slide, number):
    textbox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = textbox.text_frame
    tf.text = str(number)
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT

# ========== 幻灯片1: 封面 ==========
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
tf.text = "中国新能源汽车行业竞争格局分析"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
tf = sub_box.text_frame
tf.text = "主要玩家市占率 · 差异化策略 · 护城河评估 · 未来展望"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p.alignment = PP_ALIGN.CENTER
add_slide_number(slide, 1)

# ========== 幻灯片2-16 ==========
for slide_idx in range(2, 17):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    titles = {
        2: "核心发现摘要",
        3: "市场规模：NEV正式超越燃油车",
        4: "2024年主要玩家市占率排名",
        5: "竞争格局分层",
        6: "差异化策略矩阵",
        7: "商业模式差异化：四大流派",
        8: "护城河深度评估框架",
        9: "主要玩家护城河评级",
        10: "护城河详解：BYD vs Tesla",
        11: "护城河详解：新势力分化",
        12: "未来3年核心趋势预判",
        13: "2027年市占率预测",
        14: "竞争格局总结与投资启示",
        15: "风险提示",
        16: "谢谢观看"
    }
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = titles.get(slide_idx, "")
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    add_slide_number(slide, slide_idx)

# 保存文件
output_path = "/home/ubuntu/.openclaw-public/workspace/中国新能源汽车行业竞争格局分析.pptx"
prs.save(output_path)
print(f"PPT框架已创建，共{len(prs.slides)}页")
print(f"保存路径: {output_path}")
