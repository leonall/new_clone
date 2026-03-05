from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色方案
DARK_BLUE = RGBColor(0x1a, 0x23, 0x7e)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xd4)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xff, 0xff, 0xff)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
    p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2025年2月"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 页码
    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(len(prs.slides))
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide = add_content_slide(prs, title)
    
    # 创建表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5)).table
    
    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xe8, 0xe8, 0xe8)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = DARK_GRAY
    
    # 设置数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK_GRAY
            
            # 交替行背景色
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xfa, 0xfa, 0xfa)
    
    return slide

def add_pie_chart_slide(prs, title, categories, values, chart_title):
    """添加饼图页"""
    slide = add_content_slide(prs, title)
    
    # 创建图表数据
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(chart_title, values)
    
    # 添加图表
    x, y, cx, cy = Inches(0.5), Inches(1.8), Inches(6), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = True
    
    return slide

def add_bar_chart_slide(prs, title, categories, values, chart_title):
    """添加柱状图页"""
    slide = add_content_slide(prs, title)
    
    # 创建图表数据
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(chart_title, values)
    
    # 添加图表
    x, y, cx, cy = Inches(0.5), Inches(1.8), Inches(12), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = True
    
    return slide

# ==================== 开始创建幻灯片 ====================

# Slide 1: 标题页
add_title_slide(prs, "中国新能源汽车行业竞争格局分析", "Competitive Landscape Analysis of China NEV Industry")

# Slide 2: 执行摘要
slide = add_content_slide(prs, "执行摘要：市场格局已定，头部效应显著")
summary_text = """• 市场规模：2024年中国新能源汽车销量达1,286.6万辆，渗透率47.6%，连续十年全球第一
• 2025年预测：销量有望达1,650万辆，渗透率突破55%，增速约30%
• 竞争格局：比亚迪以34.1%市占率绝对领先，特斯拉6%位居第三
• 新势力分化：理想汽车突破50万辆成为新势力冠军，蔚来、小鹏紧随其后
• 行业趋势：价格战持续、智能化竞争加剧、出海成为新增长极"""

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5))
tf = textbox.text_frame
tf.word_wrap = True
for line in summary_text.split('\n'):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(12)

# Slide 3: 市场规模与增长趋势
slide = add_content_slide(prs, "市场规模：渗透率接近50%临界点，增长空间仍存")
market_data = [
    ["指标", "2023年", "2024年", "2025年预测"],
    ["新能源汽车销量", "949.5万辆", "1,286.6万辆", "~1,650万辆"],
    ["同比增长率", "37.9%", "35.5%", "~30%"],
    ["市场渗透率", "35.7%", "47.6%", "55%+"],
    ["全球市场份额", "~60%", "~65%", "持续提升"]
]

# 创建表格
table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.6), Inches(12.333), Inches(2)).table
for row_idx, row_data in enumerate(market_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12) if row_idx == 0 else Pt(11)
            paragraph.font.bold = (row_idx == 0)
            paragraph.font.color.rgb = DARK_GRAY
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xe8, 0xe8, 0xe8)

# 添加注释
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(2.5))
tf = note_box.text_frame
tf.word_wrap = True
notes = """关键洞察：
• 2024年新能源汽车销量突破1000万辆里程碑，渗透率接近50%关键节点
• 预计2025年国内渗透率将突破55%，部分月份可能达到60%
• 中国已连续十年保持全球最大新能源汽车市场地位，全球份额超过60%
• 出口成为新增长点，比亚迪、上汽等品牌加速全球化布局"""
for line in notes.split('\n'):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(8)

# Slide 4: 市场份额分布 - 饼图
slide = add_content_slide(prs, "市场份额：比亚迪一骑绝尘，CR5超过60%")

# 市场份额数据
categories = ['比亚迪 BYD', '吉利汽车 Geely', '特斯拉 Tesla', '五菱 SGMW', '其他 Others']
values = [34.1, 7.9, 6.0, 5.9, 46.1]

chart_data = ChartData()
chart_data.categories = categories
chart_data.add_series('2024年市场份额 (%)', values)

x, y, cx, cy = Inches(2), Inches(1.8), Inches(6), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = True

# 添加数据标签
from pptx.enum.chart import XL_DATA_LABEL_POSITION
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.show_percent = True

# Slide 5: 主要竞争对手排名
add_table_slide(prs, "主要竞争对手排名 (2024年零售销量)",
    ["排名", "厂商", "2024年销量", "市场份额", "同比增速", "主力车型"],
    [
        ["1", "比亚迪 BYD", "3,718,281辆", "34.1%", "+41.1%", "秦、宋、汉系列"],
        ["2", "吉利汽车 Geely", "862,933辆", "7.9%", "+83.2%", "极氪、银河系列"],
        ["3", "特斯拉 Tesla", "657,102辆", "6.0%", "-1.1%", "Model 3/Y"],
        ["4", "五菱 SGMW", "647,047辆", "5.9%", "+12.3%", "宏光MINI EV"],
        ["6", "理想汽车 Li Auto", "500,508辆", "4.6%", "+33.1%", "L6/L7/L8/L9"],
        ["9", "广汽埃安 GAC Aion", "~370,000辆", "3.4%", "-15%", "AION S/Y"],
        ["-", "蔚来 NIO", "221,970辆", "~2.0%", "+38.7%", "ES6/ET5/ET5T"],
        ["-", "小鹏 XPeng", "~190,000辆", "~1.7%", "+34%", "Mona M03/P7+"]
    ]
)

# Slide 6: 造车新势力对比
slide = add_content_slide(prs, "造车新势力：理想领跑，竞争格局重塑")

# 新势力对比表
new_forces_data = [
    ["品牌", "2024年销量", "同比增速", "定位", "核心技术", "2025目标"],
    ["理想 Li Auto", "500,508辆", "+33.1%", "高端家用SUV", "增程技术", "~50万辆+"],
    ["蔚来 NIO", "221,970辆", "+38.7%", "高端智能", "换电+服务", "44万辆"],
    ["小鹏 XPeng", "~190,000辆", "+34%", "智能科技", "XNGP智驾", "~30万辆"],
    ["小米 Xiaomi", "135,000辆+", "N/A", "运动轿跑", "生态整合", "30万辆"],
    ["零跑 Leapmotor", "~140,000辆", "+90%", "高性价比", "全域自研", "~25万辆"]
]

table = slide.shapes.add_table(6, 6, Inches(0.3), Inches(1.6), Inches(12.7), Inches(3.5)).table
for row_idx, row_data in enumerate(new_forces_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10) if row_idx == 0 else Pt(10)
            paragraph.font.bold = (row_idx == 0)
            paragraph.font.color.rgb = DARK_GRAY
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xe8, 0xe8, 0xe8)

# 添加分析要点
analysis_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.3), Inches(12.7), Inches(2))
tf = analysis_box.text_frame
tf.word_wrap = True
analysis = """关键变化：
• 理想汽车率先突破50万辆，成为首个达成此里程碑的新势力品牌
• 小鹏凭借Mona系列（售价11.98万起）实现销量逆袭，Mona M03月销超1.5万辆
• 小米SU7上市首年交付超13.5万辆，2025年目标30万辆
• 蔚来推出子品牌Onvo（乐道），售价下探至15-25万区间，寻求增量市场"""
for line in analysis.split('\n'):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(6)

# Slide 7: 竞争定位矩阵
slide = add_content_slide(prs, "竞争定位：价格带与技术路线分化明显")

# 创建2x2矩阵示意图
matrix_text = """┌─────────────────────────────────────────────────────────────────┐
│                          高价格带                               │
│  ┌──────────┐                           ┌──────────┐           │
│  │  蔚来    │                           │  理想    │           │
│  │ (35-60万)│                           │ (25-50万)│           │
│  └──────────┘                           └──────────┘           │
│                                                                 │
│  ┌──────────┐                           ┌──────────┐           │
│  │  特斯拉  │                           │  极氪    │           │
│  │ (23-35万)│                           │ (20-35万)│           │
│  └──────────┘                           └──────────┘           │
├─────────────────────────────────────────────────────────────────┤
│                                                                 │
│  ┌──────────┐                           ┌──────────┐           │
│  │  小鹏    │                           │  比亚迪  │           │
│  │ (12-30万)│                           │ (8-30万) │           │
│  └──────────┘                           └──────────┘           │
│                                                                 │
│  ┌──────────┐                           ┌──────────┐           │
│  │  零跑    │                           │  五菱    │           │
│  │ (5-18万) │                           │ (3-8万)  │           │
│  └──────────┘                           └──────────┘           │
│                          低价格带                               │
└─────────────────────────────────────────────────────────────────┘
         纯电BEV                              增程/混动 PHEV/EREV"""

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
tf = textbox.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = matrix_text
p.font.name = 'Courier New'
p.font.size = Pt(10)
p.font.color.rgb = DARK_GRAY

# Slide 8: 护城河分析
slide = add_content_slide(prs, "护城河评估：规模、技术与生态构建竞争壁垒")

moat_table = [
    ["品牌", "规模优势", "技术护城河", "品牌/服务", "综合评级"],
    ["比亚迪", "●●● 全球最大", "●●● 刀片电池+DM-i", "●●○ 国民品牌", "★★★★★"],
    ["特斯拉", "●●○ 全球第六", "●●● FSD+超充网络", "●●● 科技标杆", "★★★★★"],
    ["理想", "●●○ 新势力第一", "●●○ 增程方案", "●●● 家庭定位", "★★★★☆"],
    ["蔚来", "●○○ 规模较小", "●●● 换电网络", "●●● 服务标杆", "★★★★☆"],
    ["小鹏", "●○○ 快速追赶", "●●● 智能驾驶", "●●○ 科技年轻", "★★★☆☆"],
    ["小米", "●○○ 起步阶段", "●●○ 生态整合", "●●○ 流量入口", "★★★☆☆"]
]

table = slide.shapes.add_table(7, 5, Inches(0.3), Inches(1.6), Inches(12.7), Inches(4.5)).table
for row_idx, row_data in enumerate(moat_table):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11) if row_idx == 0 else Pt(10)
            paragraph.font.bold = (row_idx == 0)
            paragraph.font.color.rgb = DARK_GRAY
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xe8, 0xe8, 0xe8)

# Slide 9: 行业趋势与挑战
slide = add_content_slide(prs, "行业趋势：智能化竞赛加速，出海成为必选项")
trends = """🔹 价格战持续白热化
   • 2024年新能源汽车平均售价下降约10-15%
   • 比亚迪、特斯拉多次调价，倒逼行业降本
   
🔹 智能化成为核心差异化
   • 城市NOA（导航辅助驾驶）成为标配竞争点
   • 端到端大模型上车，智驾能力快速迭代
   • 座舱智能化向AI Agent演进
   
🔹 出海战略提速
   • 2024年中国新能源汽车出口超200万辆
   • 比亚迪、奇瑞、上汽在东南亚、拉美市场份额提升
   • 欧洲市场面临关税壁垒，本土化建厂成为趋势
   
🔹 产业链垂直整合深化
   • 电池、电机、电控等核心部件自研比例提升
   • 智能驾驶芯片国产替代加速（华为、地平线）"""

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
tf = textbox.text_frame
tf.word_wrap = True
for line in trends.split('\n'):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(6)

# Slide 10: 战略建议
slide = add_content_slide(prs, "战略建议：存量博弈时代的生存法则")

recommendations = """对于现有参与者：

1. 规模效应是生存基础
   • 年销量30万辆是新势力盈亏平衡线，50万辆是安全边际
   • 理想、蔚来已突破，小鹏、零跑加速追赶

2. 差异化定位避免同质化竞争
   • 理想：家庭场景 + 增程路线
   • 蔚来：高端服务 + 换电生态
   • 小鹏：智能驾驶 + 科技普惠

3. 成本控制与技术创新并重
   • 电池成本占整车30-40%，垂直整合能力关键
   • 智驾能力决定溢价空间，但投入巨大

4. 海外市场是第二增长曲线
   • 东南亚、中东、拉美市场机会窗口期
   • 欧洲市场需应对贸易壁垒，本土化生产是长期方案

⚠️ 风险提示：行业整合加速，2025-2026年可能迎来洗牌期"""

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
tf = textbox.text_frame
tf.word_wrap = True
for line in recommendations.split('\n'):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(6)

# Slide 11: 数据来源与免责声明
slide = add_content_slide(prs, "数据来源与免责声明")
disclaimer = """数据来源：
• 中国汽车工业协会 (CAAM)
• 乘用车市场信息联席会 (CPCA)
• 各公司财报及月度交付报告
• CnEVPost, CarNewsChina 等行业媒体

⚠️ 免责声明：
本报告仅供行业研究和学习参考，不构成任何投资建议。
市场数据基于公开信息整理，可能存在时效性差异。
投资有风险，决策需谨慎。

报告制作时间：2025年2月"""

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
tf = textbox.text_frame
tf.word_wrap = True
for line in disclaimer.split('\n'):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(10)

# 保存PPT
output_path = "/home/ubuntu/.openclaw-public/workspace/中国新能源汽车行业竞争格局分析_2025.pptx"
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
