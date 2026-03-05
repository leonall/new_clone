from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题
NAVY = RGBColor(0x1A, 0x23, 0x7E)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)

# 辅助函数：添加标题页
def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.alignment = PP_ALIGN.CENTER
    
    add_slide_number(slide, len(prs.slides))
    return slide

# 辅助函数：添加内容页
def add_content_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    add_slide_number(slide, len(prs.slides))
    return slide

# 添加页码
def add_slide_number(slide, number):
    textbox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = textbox.text_frame
    tf.text = str(number)
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT

# ========== 开始创建幻灯片 ==========

# Slide 1: 封面
add_title_slide(prs, 
    "中国新能源汽车行业竞争格局分析",
    "主要玩家市占率 · 差异化策略 · 护城河评估 · 未来展望")

# Slide 2: 核心发现摘要
slide = add_content_slide(prs, "核心发现摘要")
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

findings = [
    ("市场格局", "\"一超多强\"格局稳固，BYD以34.1%市占率遥遥领先，CR3达48%"),
    ("竞争分层", "Tier 1: BYD绝对龙头 | Tier 2: Geely/Tesla/五菱 | Tier 3: 新势力分化"),
    ("差异化策略", "技术路线分化：纯电vs增程；智能化成为核心战场"),
    ("护城河评估", "BYD垂直整合最强，Tesla技术+品牌双驱动，新势力各有特色"),
    ("未来展望", "2025-2027年集中度将进一步提升，智能化成为决胜关键，预计CR3达55%+"),
    ("投资启示", "短期关注盈利能力，中期关注智能化进展，长期看好出海布局")
]

for i, (title, desc) in enumerate(findings):
    if i > 0:
        tf.add_paragraph()
    p = tf.add_paragraph()
    p.text = f"{title}："
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_before = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    p2.space_after = Pt(6)
    p2.level = 1

# Slide 3: 市场规模
slide = add_content_slide(prs, "市场规模：NEV正式超越燃油车")

metrics = [
    ("1,088万辆", "2024年NEV销量", Inches(0.7), Inches(1.6)),
    ("50%+", "NEV渗透率", Inches(4.0), Inches(1.6)),
    ("5-8%", "2025-2027 CAGR", Inches(7.3), Inches(1.6)),
    ("34.1%", "BYD市占率", Inches(10.6), Inches(1.6))
]

for value, label, left, top in metrics:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(2)
    
    val_box = slide.shapes.add_textbox(left, top + Inches(0.3), Inches(2.8), Inches(0.8))
    tf = val_box.text_frame
    tf.text = value
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    label_box = slide.shapes.add_textbox(left, top + Inches(1.1), Inches(2.8), Inches(0.5))
    tf = label_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.6))
tf = note_box.text_frame
tf.text = "数据来源：CPCA（中国汽车流通协会）| NEV包含BEV纯电动+PHEV插电混动+EREV增程"
p = tf.paragraphs[0]
p.font.size = Pt(10)
p.font.color.rgb = DARK_GRAY
p.font.italic = True

# Slide 4: 市占率排名表
slide = add_content_slide(prs, "2024年主要玩家市占率排名")

rows, cols = 12, 5
table = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(1.5), Inches(12), Inches(5.2)).table

table.columns[0].width = Inches(0.8)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.2)
table.columns[4].width = Inches(2.0)

headers = ["排名", "品牌", "市占率", "销量(万辆)", "同比增长"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(12)
    paragraph.font.bold = True
    paragraph.font.color.rgb = WHITE
    paragraph.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

data = [
    ("1", "BYD", "34.1%", "427", "+41%"),
    ("2", "Geely", "7.9%", "86", "+45%"),
    ("3", "Tesla", "6.0%", "66", "-3%"),
    ("4", "五菱", "5.9%", "65", "-5%"),
    ("5", "长安汽车", "5.2%", "57", "+38%"),
    ("6", "理想汽车", "4.6%", "50", "+33%"),
    ("7", "赛力斯(AITO)", "3.8%", "41", "+182%"),
    ("8", "零跑汽车", "3.6%", "39", "+104%"),
    ("9", "广汽埃安", "3.4%", "37", "-21%"),
    ("10", "小鹏汽车", "3.1%", "34", "+34%"),
    ("11", "蔚来汽车", "2.8%", "31", "+39%")
]

for row_idx, row_data in enumerate(data, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(11)
        paragraph.alignment = PP_ALIGN.CENTER if col_idx != 1 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        
        if col_idx == 4 and value.startswith("-"):
            paragraph.font.color.rgb = RED

# Slide 5: 竞争格局分层
slide = add_content_slide(prs, "竞争格局分层：\"一超多强\"态势明显")

tiers = [
    ("Tier 1 绝对龙头", "BYD", "销量427万辆，市占率34.1%\n全价格带覆盖，垂直整合优势", 
     RGBColor(0x00, 0x66, 0xCC), Inches(0.7)),
    ("Tier 2 传统+外资", "Geely / Tesla / 五菱", "市占率5-8%区间\n多品牌/纯电标杆/极致性价比", 
     RGBColor(0x66, 0x99, 0xCC), Inches(4.5)),
    ("Tier 3 新势力", "理想/问界/零跑/小鹏/蔚来", "市占率2-5%，快速分化\n各有差异化定位", 
     RGBColor(0x99, 0xCC, 0xFF), Inches(8.3))
]

for tier_name, players, desc, color, left in tiers:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5.2))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    
    tier_box = slide.shapes.add_textbox(left, Inches(1.7), Inches(3.8), Inches(0.6))
    tf = tier_box.text_frame
    tf.text = tier_name
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    player_box = slide.shapes.add_textbox(left, Inches(2.4), Inches(3.8), Inches(0.8))
    tf = player_box.text_frame
    tf.text = players
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    desc_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(3.3), Inches(3.4), Inches(3.2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.text = desc
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.line_spacing = 1.3

# Slide 6: 差异化策略矩阵
slide = add_content_slide(prs, "差异化策略矩阵：技术路线分化")

rows, cols = 8, 4
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(12.4), Inches(5.6)).table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(3.2)
table.columns[2].width = Inches(3.8)
table.columns[3].width = Inches(3.4)

headers = ["品牌", "动力路线", "智能化策略", "目标客群"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

strategy_data = [
    ("BYD", "纯电+插混并重", "自研芯片+DiLink系统", "大众市场"),
    ("Tesla", "纯电", "纯视觉FSD方案", "科技先锋"),
    ("理想", "增程为主", "家庭智能座舱", "家庭用户"),
    ("问界", "增程+纯电", "华为ADS 3.0智驾", "商务/科技"),
    ("小鹏", "纯电", "城市NGP领航", "年轻科技"),
    ("蔚来", "纯电+换电", "NOMI+NOP辅助驾驶", "高端商务"),
    ("零跑", "纯电+增程", "自研芯片降本", "性价比")
]

for row_idx, row_data in enumerate(strategy_data, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT if col_idx > 0 else PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

print("幻灯片1-6创建完成")

# Slide 7: 商业模式对比
slide = add_content_slide(prs, "商业模式差异化：四大流派")

models = [
    ("垂直整合派", "BYD", 
     "• 自研自产电池、电机、电控\n• 自研IGBT/SiC芯片\n• 成本优势：较竞品低15-20%",
     Inches(0.6), Inches(1.5)),
    ("生态赋能派", "问界/华为",
     "• 华为提供智驾+座舱\n• 共享华为门店渠道\n• 轻资产快速起量",
     Inches(6.8), Inches(1.5)),
    ("用户运营派", "蔚来",
     "• 换电网络（3400+座）\n• 用户社区+生活方式\n• 高NPS、高复购率",
     Inches(0.6), Inches(4.4)),
    ("极致效率派", "零跑/理想",
     "• 聚焦爆款单品\n• 平台化开发低SKU\n• 高周转、快迭代",
     Inches(6.8), Inches(4.4))
]

for model_name, example, features, left, top in models:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.0), Inches(2.6))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(2)
    
    name_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(5.7), Inches(0.5))
    tf = name_box.text_frame
    tf.text = model_name
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    ex_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.6), Inches(5.7), Inches(0.4))
    tf = ex_box.text_frame
    tf.text = f"代表：{example}"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    feat_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(1.0), Inches(5.7), Inches(1.5))
    tf = feat_box.text_frame
    tf.word_wrap = True
    tf.text = features
    for p in tf.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

# Slide 8: 护城河评估框架
slide = add_content_slide(prs, "护城河深度评估框架")

moat_dims = [
    ("规模经济", "成本曲线优势\n• 产能利用率\n• 采购议价权\n• 制造效率", Inches(0.7), Inches(1.6)),
    ("技术壁垒", "自研核心能力\n• 电池技术\n• 智能驾驶\n• 芯片/电控", Inches(4.5), Inches(1.6)),
    ("品牌溢价", "定价权与用户忠诚度\n• 品牌认知度\n• 用户NPS\n• 定价能力", Inches(8.3), Inches(1.6)),
    ("渠道网络", "触达与服务覆盖\n• 销售网点数\n• 下沉市场覆盖\n• 服务便利性", Inches(0.7), Inches(4.5)),
    ("生态系统", "软硬一体化\n• 软件服务收入\n• 充电/换电网络\n• 数据闭环", Inches(4.5), Inches(4.5))
]

for dim_name, desc, left, top in moat_dims:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4.0), Inches(2.6))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = NAVY
    card.line.width = Pt(1.5)
    
    name_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(3.7), Inches(0.5))
    tf = name_box.text_frame
    tf.text = dim_name
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    desc_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.65), Inches(3.7), Inches(1.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.text = desc
    for p in tf.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY

# Slide 9: 护城河评级表
slide = add_content_slide(prs, "主要玩家护城河评级")

rows, cols = 8, 7
table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.6)).table

for i in range(cols):
    table.columns[i].width = Inches(1.8)
table.columns[0].width = Inches(1.5)

headers = ["品牌", "规模经济", "技术壁垒", "品牌溢价", "渠道网络", "生态系统", "综合评级"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

ratings_data = [
    ("BYD", "●●●", "●●●", "●●○", "●●●", "●●○", "强"),
    ("Tesla", "●●○", "●●●", "●●●", "●●○", "●●●", "强"),
    ("理想", "●●○", "●●○", "●●●", "●●○", "●●○", "中"),
    ("问界", "●●○", "●●●", "●●○", "●●○", "●●○", "中"),
    ("小鹏", "●○○", "●●●", "●○○", "●○○", "●●○", "中"),
    ("蔚来", "●○○", "●●○", "●●○", "●○○", "●●●", "中"),
    ("零跑", "●○○", "●●○", "●○○", "●●○", "●○○", "弱-中")
]

rating_colors = {
    "●●●": GREEN,
    "●●○": ORANGE,
    "●○○": RED,
    "强": GREEN,
    "中": ORANGE,
    "弱-中": RED
}

for row_idx, row_data in enumerate(ratings_data, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if col_idx > 0 and value in rating_colors:
            p.font.color.rgb = rating_colors[value]
        
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

legend_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
tf = legend_box.text_frame
tf.text = "评级说明：●●● 强  |  ●●○ 中  |  ●○○ 弱"
p = tf.paragraphs[0]
p.font.size = Pt(10)
p.font.color.rgb = DARK_GRAY

print("幻灯片7-9创建完成")

# Slide 10: BYD vs Tesla
slide = add_content_slide(prs, "护城河详解：BYD vs Tesla")

byd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5))
byd_box.fill.solid()
byd_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xF8)
byd_box.line.color.rgb = ACCENT_BLUE
byd_box.line.width = Pt(2)

byd_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.5))
tf = byd_title.text_frame
tf.text = "BYD：垂直整合之王"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

byd_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.6), Inches(4.5))
tf = byd_content.text_frame
tf.word_wrap = True
content = """优势：
• 刀片电池技术领先，安全口碑
• 自研IGBT/SiC芯片，供应链安全
• 成本优势：较竞品低15-20%
• 全价格带覆盖能力

风险：
• 品牌向上难度大
• 智能化相对落后
• 海外扩张挑战"""
tf.text = content
for p in tf.paragraphs:
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

tesla_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5))
tesla_box.fill.solid()
tesla_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
tesla_box.line.color.rgb = RED
tesla_box.line.width = Pt(2)

tesla_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5))
tf = tesla_title.text_frame
tf.text = "Tesla：技术+品牌双驱动"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RED

tesla_content = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.6), Inches(4.5))
tf = tesla_content.text_frame
tf.word_wrap = True
content = """优势：
• FSD技术领先，数据闭环
• 品牌溢价：售价高于同级20%+
• 超充网络生态
• 制造创新（一体化压铸）

风险：
• 产品迭代放缓
• 中国本土化不足
• 价格战压力"""
tf.text = content
for p in tf.paragraphs:
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

# Slide 11: 新势力详解
slide = add_content_slide(prs, "护城河详解：新势力分化")

startups = [
    ("理想", "家庭场景定义者", 
     "• 精准定位增程家庭SUV\n• 高毛利率20%+，已盈利\n• 产品定义能力强",
     "• 纯电转型压力\n• 依赖单一细分市场",
     Inches(0.4)),
    ("问界", "华为生态加持",
     "• 华为ADS智驾国内领先\n• 华为门店渠道覆盖\n• 品牌背书强",
     "• 过度依赖华为\n• 品牌独立性弱",
     Inches(3.4)),
    ("小鹏", "智能驾驶标签",
     "• XNGP城市领航领先\n• 技术型人才密度高\n• 智驾数据积累",
     "• 品牌定位模糊\n• 盈利压力大",
     Inches(6.4)),
    ("蔚来", "用户运营标杆",
     "• 换电网络壁垒（3400+座）\n• 用户社区运营领先\n• 高端品牌形象",
     "• 重资产模式\n• 持续亏损\n• 资金压力大",
     Inches(9.4))
]

for name, tagline, pros, cons, left in startups:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(3.2), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = NAVY
    card.line.width = Pt(1)
    
    name_tb = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.5), Inches(3.0), Inches(0.5))
    tf = name_tb.text_frame
    tf.text = name
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    tag_tb = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.95), Inches(3.0), Inches(0.4))
    tf = tag_tb.text_frame
    tf.text = tagline
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_BLUE
    p.font.italic = True
    
    pros_tb = slide.shapes.add_textbox(left + Inches(0.1), Inches(2.4), Inches(3.0), Inches(1.8))
    tf = pros_tb.text_frame
    tf.word_wrap = True
    tf.text = pros
    for p in tf.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_GRAY
    
    cons_tb = slide.shapes.add_textbox(left + Inches(0.1), Inches(4.3), Inches(3.0), Inches(1.8))
    tf = cons_tb.text_frame
    tf.word_wrap = True
    tf.text = "风险：\n" + cons
    for p in tf.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RED

print("幻灯片10-11创建完成")

# Slide 12: 未来趋势
slide = add_content_slide(prs, "未来3年核心趋势预判（2025-2027）")

trends = [
    ("趋势一", "价格战趋缓，价值战开启", 
     "• 2024-2025年价格战触底\n• 2025H2转向智能化、服务差异化\n• 行业毛利率缓慢修复"),
    ("趋势二", "智能化成为核心战场", 
     "• 城市NOA成为标配（2025年底）\n• 端到端大模型重塑智驾格局\n• 座舱AI Agent差异化竞争"),
    ("趋势三", "市场集中度提升", 
     "• CR3有望从48%提升至55%+\n• 尾部品牌加速出清\n• 新势力分化：2-3家胜出"),
    ("趋势四", "出海成为第二增长曲线", 
     "• 2025年出口占比预计达20%\n• 欧洲、东南亚、拉美为重点\n• 本地化建厂规避关税壁垒")
]

for i, (num, title, desc) in enumerate(trends):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.5 + row * 2.9)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.1), Inches(2.7))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(2)
    
    num_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.1), top + Inches(0.1), Inches(0.9), Inches(0.45))
    num_bg.fill.solid()
    num_bg.fill.fore_color.rgb = NAVY
    num_bg.line.fill.background()
    
    num_tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(0.8), Inches(0.4))
    tf = num_tb.text_frame
    tf.text = num
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    title_tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.7), Inches(5.8), Inches(0.5))
    tf = title_tb.text_frame
    tf
