from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1A, 0x23, 0x7E)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)

def add_slide_number(slide, number):
    textbox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = textbox.text_frame
    tf.text = str(number)
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_title_bar(slide, title_text):
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

# Slide 1: 封面
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
tf.text = "中国新能源汽车行业竞争格局分析"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
tf = sub_box.text_frame
tf.text = "主要玩家市占率 · 差异化策略 · 护城河评估 · 未来展望"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p.alignment = PP_ALIGN.CENTER
add_slide_number(slide, 1)

# Slide 2: 核心发现
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "核心发现摘要")
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True
findings = [
    "【市场格局】\"一超多强\"格局稳固，BYD以34.1%市占率遥遥领先，CR3达48%",
    "【竞争分层】Tier 1: BYD绝对龙头 | Tier 2: Geely/Tesla/五菱 | Tier 3: 新势力分化",
    "【差异化策略】技术路线分化：纯电vs增程；智能化成为核心战场",
    "【护城河评估】BYD垂直整合最强，Tesla技术+品牌双驱动，新势力各有特色",
    "【未来展望】2025-2027年集中度将进一步提升，智能化成为决胜关键",
    "【投资启示】短期关注盈利能力，中期关注智能化进展，长期看好出海布局"
]
for i, text in enumerate(findings):
    if i > 0:
        tf.add_paragraph()
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(16)
add_slide_number(slide, 2)

# Slide 3-16: 其他页面
page_titles = [
    "市场规模：NEV正式超越燃油车",
    "2024年主要玩家市占率排名",
    "竞争格局分层：\"一超多强\"态势明显",
    "差异化策略矩阵：技术路线分化",
    "商业模式差异化：四大流派",
    "护城河深度评估框架",
    "主要玩家护城河评级",
    "护城河详解：BYD vs Tesla",
    "护城河详解：新势力分化",
    "未来3年核心趋势预判（2025-2027）",
    "2027年市占率预测与情景分析",
    "竞争格局总结与投资启示",
    "风险提示"
]

for idx, title in enumerate(page_titles, 3):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    
    # 添加内容占位
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = content.text_frame
    tf.word_wrap = True
    
    if idx == 3:  # 市场规模
        tf.text = """关键数据：
• 2024年NEV销量：1,088万辆（同比+35%）
• NEV渗透率：>50%（首次超越燃油车）
• 预计2025-2027年CAGR：5-8%
• BYD市占率：34.1%（遥遥领先）

数据来源：CPCA中国汽车流通协会"""
    elif idx == 4:  # 市占率
        tf.text = """2024年TOP 11品牌市占率：

1. BYD        34.1%  |  销量427万辆  |  +41%
2. Geely       7.9%  |  销量86万辆   |  +45%
3. Tesla       6.0%  |  销量66万辆   |  -3%
4. 五菱        5.9%  |  销量65万辆   |  -5%
5. 长安        5.2%  |  销量57万辆   |  +38%
6. 理想        4.6%  |  销量50万辆   |  +33%
7. 问界(AITO)  3.8%  |  销量41万辆   |  +182%
8. 零跑        3.6%  |  销量39万辆   |  +104%
9. 埃安        3.4%  |  销量37万辆   |  -21%
10. 小鹏       3.1%  |  销量34万辆   |  +34%
11. 蔚来       2.8%  |  销量31万辆   |  +39%"""
    elif idx == 5:  # 竞争分层
        tf.text = """Tier 1 绝对龙头：BYD
• 市占率34.1%，销量427万辆
• 全价格带覆盖（8万-80万）

Tier 2 传统+外资：Geely/Tesla/五菱
• 市占率5-8%区间
• 多品牌/纯电标杆/极致性价比

Tier 3 新势力：理想/问界/零跑/小鹏/蔚来
• 市占率2-5%，快速分化
• 各有差异化定位"""
    elif idx == 6:  # 差异化策略
        tf.text = """技术路线分化：

BYD    | 纯电+插混并重 | 自研芯片+DiLink   | 大众市场
Tesla  | 纯电          | 纯视觉FSD        | 科技先锋
理想   | 增程为主      | 家庭智能座舱     | 家庭用户
问界   | 增程+纯电     | 华为ADS 3.0      | 商务/科技
小鹏   | 纯电          | 城市NGP领航      | 年轻科技
蔚来   | 纯电+换电     | NOMI+NOP         | 高端商务
零跑   | 纯电+增程     | 自研芯片降本     | 性价比"""
    elif idx == 7:  # 商业模式
        tf.text = """四大商业模式流派：

【垂直整合派】BYD
• 自研自产电池/电机/电控/芯片
• 成本优势：较竞品低15-20%

【生态赋能派】问界/华为
• 华为提供智驾+座舱+渠道
• 轻资产运营，快速起量

【用户运营派】蔚来
• 换电网络(3400+座)+用户社区
• 高NPS、高复购率

【极致效率派】零跑/理想
• 聚焦爆款单品+平台化开发
• 低SKU、高周转"""
    elif idx == 8:  # 护城河框架
        tf.text = """五大护城河维度：

1. 规模经济：成本曲线优势
   • 产能利用率/采购议价权/制造效率

2. 技术壁垒：自研核心能力
   • 电池技术/智能驾驶/芯片电控

3. 品牌溢价：定价权与用户忠诚度
   • 品牌认知度/NPS/定价能力

4. 渠道网络：触达与服务覆盖
   • 销售网点/下沉覆盖/服务便利

5. 生态系统：软硬一体化
   • 软件服务/充电换电/数据闭环"""
    elif idx == 9:  # 护城河评级
        tf.text = """护城河评级（●●●强 ●●○中 ●○○弱）：

品牌    规模经济  技术壁垒  品牌溢价  渠道网络  生态系统  综合
BYD      ●●●      ●●●      ●●○      ●●●      ●●○      强
Tesla    ●●○      ●●●      ●●●      ●●○      ●●●      强
理想     ●●○      ●●○      ●●●      ●●○      ●●○      中
问界     ●●○      ●●●      ●●○      ●●○      ●●○      中
小鹏     ●○○      ●●●      ●○○      ●○○      ●●○      中
蔚来     ●○○      ●●○      ●●○      ●○○      ●●●      中
零跑     ●○○      ●●○      ●○○      ●●○      ●○○     弱-中"""
    elif idx == 10:  # BYD vs Tesla
        tf.text = """【BYD：垂直整合之王】
优势：刀片电池技术领先/自研芯片/成本优势15-20%/全价格带覆盖
风险：品牌向上难/智能化落后/海外扩张挑战

【Tesla：技术+品牌双驱动】
优势：FSD技术领先/品牌溢价20%+/超充网络/制造创新
风险：产品迭代放缓/中国本土化不足/价格战压力"""
    elif idx == 11:  # 新势力
        tf.text = """新势力护城河分化：

【理想】家庭场景定义者
优势：增程家庭SUV定位/毛利率20%+/已盈利
风险：纯电转型压力/依赖单一细分市场

【问界】华为生态加持
优势：华为ADS智驾领先/渠道覆盖/品牌背书
风险：过度依赖华为/品牌独立性弱

【小鹏】智能驾驶标签
优势：XNGP城市领航/技术人才密度
风险：品牌定位模糊/盈利压力大

【蔚来】用户运营标杆
优势：换电网络(3400+座)/用户社区/高端品牌
风险：重资产模式/持续亏损/资金压力"""
    elif idx == 12:  # 未来趋势
        tf.text = """2025-2027年核心趋势：

趋势一：价格战趋缓，价值战开启
• 2024-2025年价格战触底
• 2025H2转向智能化/服务差异化

趋势二：智能化成为核心战场
• 城市NOA成为标配（2025年底）
• 端到端大模型重塑智驾格局

趋势三：市场集中度提升
• CR3从48%提升至55%+
• 尾部品牌加速出清

趋势四：出海成为第二增长曲线
• 2025年出口占比预计达20%
• 欧洲、东南亚、拉美为重点"""
    elif idx == 13:  # 预测
        tf.text = """2027年市占率预测：

品牌      2024市占率    2027预测      趋势
BYD        34.1%       28-32%        ↓
Tesla       6.0%        5-7%         ↓
理想        4.6%        4-6%          →
问界        3.8%        4-6%          ↑
零跑        3.6%        4-5%          ↑
小鹏        3.1%        3-4%          →
蔚来        2.8%        2-3%          ↓

情景概率：基准50% | 乐观25% | 悲观25%"""
    elif idx == 14:  # 总结
        tf.text = """竞争格局总结与投资启示：

【格局总结】
• BYD：短期无忧，长期需关注智能化与品牌向上
• Tesla：技术护城河仍在，但中国市场挑战加剧
• 新势力：理想最稳健，问界增速最快，零跑黑马潜力
• 行业：从\"拼规模\"进入\"拼效率+拼技术\"阶段

【投资启示】
短期 → 关注盈利能力
中期 → 关注智能化进展
长期 → 看好出海布局"""
    else:  # 风险提示
        tf.text = """主要风险因素：

【技术风险】
• 智能驾驶技术迭代不及预期
• 电池技术路线变化
• 芯片供应受限

【市场风险】
• 宏观经济下行影响消费需求
• 价格战持续导致行业性亏损

【国际风险】
• 国际贸易摩擦升级
• 关税壁垒增加
• 地缘政治影响出海

【竞争风险】
• 新进入者颠覆
• 技术路线押注错误

⚠️ 免责声明：本分析仅供参考和学习，不构成投资建议。"""
    
    for p in tf.paragraphs:
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
    
    add_slide_number(slide, idx)

# 保存
output_path = "/home/ubuntu/.openclaw-public/workspace/中国新能源汽车行业竞争格局分析.pptx"
prs.save(output_path)
print(f"PPT已保存: {output_path}")
print(f"共{len(prs.slides)}页幻灯片")
