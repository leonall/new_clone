"""
中国新能源汽车行业竞争格局分析 - Enhanced with embedded charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
import copy

# ─── Color palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x23, 0x7E)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY  = RGBColor(0x75, 0x75, 0x75)
ACCENT    = RGBColor(0x00, 0x7A, 0xCC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
ORANGE    = RGBColor(0xEF, 0x6C, 0x00)
RED       = RGBColor(0xC6, 0x28, 0x28)
GOLD      = RGBColor(0xF5, 0x9E, 0x0E)

# Brand colors for chart series
CHART_COLORS = [
    RGBColor(0x00, 0x7A, 0xCC),   # blue - BYD
    RGBColor(0xEF, 0x6C, 0x00),   # orange - Geely
    RGBColor(0xC6, 0x28, 0x28),   # red - Tesla
    RGBColor(0x2E, 0x7D, 0x32),   # green - SGMW
    RGBColor(0x6A, 0x1B, 0x9A),   # purple - Changan
    RGBColor(0x00, 0x83, 0x8A),   # teal - Li Auto
    RGBColor(0xF5, 0x9E, 0x0E),   # gold - AITO
    RGBColor(0x37, 0x47, 0x4F),   # slate - Leapmotor
    RGBColor(0x55, 0x8B, 0x2F),   # olive - GAC Aion
    RGBColor(0xAD, 0x14, 0x57),   # pink - Xpeng
    RGBColor(0x1B, 0x5E, 0x20),   # dark green - NIO
]

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def title_bar(slide, text, subtitle=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, Inches(1.1))
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY
    rect.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15),
                                   Inches(12.333), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle
        p2.font.size = Pt(13); p2.font.color.rgb = RGBColor(0xBB,0xCC,0xDD)

def page_num(slide, n):
    tb = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = str(n); p.font.size = Pt(11); p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT

def source_note(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(12), Inches(0.3))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(9); p.font.color.rgb = MID_GRAY

def add_textbox(slide, left, top, width, height, text, size=13,
                bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color or DARK_GRAY
    p.alignment = align
    return tb

def add_table(slide, left, top, width, height, data, header_color=NAVY,
              col_widths=None):
    """data = list of rows; first row = header"""
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols,
        Inches(left), Inches(top), Inches(width), Inches(height)).table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(width * w / total)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
            else:
                p.font.color.rgb = DARK_GRAY
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5,0xF7,0xFA)
    return tbl

def set_series_color(chart, series_idx, rgb):
    """Set fill color for a chart series."""
    series = chart.series[series_idx]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = rgb

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Cover
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Accent strip
strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(5.8), prs.slide_width, Inches(0.08))
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_textbox(slide, 1.5, 1.8, 10, 1.8,
            "中国新能源汽车行业竞争格局分析",
            size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.8, 10, 0.6,
            "主要玩家市占率  ·  差异化策略  ·  护城河评估  ·  未来展望",
            size=18, color=RGBColor(0xBB,0xCC,0xEE), align=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.6, 10, 0.5,
            "数据截至 2024年12月  |  数据来源：CPCA / 各公司年报",
            size=13, color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)
page_num(slide, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Executive Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "核心发现摘要", "6大关键洞察")
page_num(slide, 2)

bullets = [
    ("🏆 市场格局",   "\"一超多强\"稳固，BYD市占率34.1%，CR3达48%，市场集中度持续提升"),
    ("🔢 竞争分层",   "Tier 1 BYD独大 → Tier 2 Geely/Tesla/五菱 → Tier 3 新势力分化"),
    ("⚙️ 差异化路线", "技术路线分化：纯电 vs 增程；2025年智能化(NOA/端到端)成核心战场"),
    ("🏰 护城河评估", "BYD垂直整合护城河最深；Tesla技术+品牌双轮驱动；理想盈利率领跑新势力"),
    ("📈 未来展望",   "2025-2027年集中度再提升：BYD 28-32%；出海成增量，东南亚/欧洲为主"),
    ("💡 投资启示",   "短期关注盈利可持续性，中期关注智能化进展，长期看好具备出海能力的龙头"),
]
colors = [ACCENT, NAVY, GREEN, ORANGE, RED, RGBColor(0x6A,0x1B,0x9A)]

for i, (tag, body) in enumerate(bullets):
    row, col = divmod(i, 2)
    lx = 0.4 + col * 6.5; ty = 1.35 + row * 1.75
    # card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(ty), Inches(6.3), Inches(1.55))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xDD,0xE2,0xEE)
    # colored left bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(ty), Inches(0.08), Inches(1.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = colors[i]
    bar.line.fill.background()
    add_textbox(slide, lx+0.15, ty+0.08, 6.0, 0.4, tag,
                size=13, bold=True, color=colors[i])
    add_textbox(slide, lx+0.15, ty+0.48, 5.9, 0.95, body,
                size=11.5, color=DARK_GRAY, wrap=True)

source_note(slide, "来源：CPCA 2024、各公司年报及公开投资者文件")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Market Size (with bar chart)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "市场规模：NEV渗透率突破50%，正式超越燃油车")
page_num(slide, 3)

# KPI boxes
kpis = [
    ("1,088万辆", "2024年NEV零售销量"),
    (">50%",      "NEV渗透率（首超燃油车）"),
    ("+35%",      "2024年销量同比增速"),
    ("5-8%",      "2025-2027E CAGR"),
]
for i, (val, lbl) in enumerate(kpis):
    lx = 0.4 + i * 3.1
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.25), Inches(2.9), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    add_textbox(slide, lx+0.1, 1.3, 2.7, 0.6, val,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, lx+0.1, 1.85, 2.7, 0.5, lbl,
                size=11, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)

# Penetration rate bar chart (2020-2024)
cd = CategoryChartData()
cd.categories = ["2020", "2021", "2022", "2023", "2024"]
cd.add_series("NEV渗透率(%)", (5.4, 13.4, 25.6, 35.7, 51.2))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(2.6), Inches(7.0), Inches(4.4), cd
).chart
chart.has_legend = False
chart.has_title = True
chart.chart_title.text_frame.text = "中国NEV渗透率趋势（2020-2024）"
chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
chart.chart_title.text_frame.paragraphs[0].font.bold = True
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.number_format = '0.0"%"'
chart.plots[0].data_labels.font.size = Pt(10)
set_series_color(chart, 0, ACCENT)

# Volume bar chart
cd2 = CategoryChartData()
cd2.categories = ["2020", "2021", "2022", "2023", "2024"]
cd2.add_series("销量(万辆)", (136, 352, 649, 773, 1088))
chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.7), Inches(2.6), Inches(5.3), Inches(4.4), cd2
).chart
chart2.has_legend = False
chart2.has_title = True
chart2.chart_title.text_frame.text = "NEV年销量趋势（万辆）"
chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
chart2.chart_title.text_frame.paragraphs[0].font.bold = True
chart2.plots[0].has_data_labels = True
chart2.plots[0].data_labels.font.size = Pt(10)
set_series_color(chart2, 0, GREEN)

source_note(slide, "来源：CPCA（乘联会）2024年数据；2025-2027年为分析师预测")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Market Share (PIE chart)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "2024年市占率：BYD以34.1%遥遥领先，CR3达48%")
page_num(slide, 4)

brands     = ["BYD","Geely","Tesla","SGMW","Changan","Li Auto","AITO","Leapmotor","GAC Aion","Xpeng","NIO","Others"]
shares     = [34.1, 7.9, 6.0, 5.9, 5.2, 4.6, 3.8, 3.6, 3.4, 3.1, 2.8, 19.6]
volumes    = [427, 86, 66, 65, 57, 50, 41, 39, 37, 34, 31, None]
yoy        = ["+41%","+45%","-3%","-5%","+38%","+33%","+182%","+104%","-21%","+34%","+39%","—"]

cd = CategoryChartData()
cd.categories = brands
cd.add_series("市占率", shares)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(0.3), Inches(1.2), Inches(6.5), Inches(5.9), cd
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = True
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_percentage = True
chart.plots[0].data_labels.show_category_name = False
chart.plots[0].data_labels.font.size = Pt(9)

# Colour the first 11 slices
for idx, rgb in enumerate(CHART_COLORS):
    try:
        pt = chart.plots[0].series[0].points[idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = rgb
    except Exception:
        pass

# Table on the right
header = ["品牌", "市占率", "销量(万辆)", "同比"]
rows = [header]
for i in range(len(brands)-1):
    rows.append([brands[i], f"{shares[i]}%",
                 str(volumes[i]) if volumes[i] else "—", yoy[i]])
add_table(slide, 7.1, 1.2, 6.0, 5.9, rows,
          col_widths=[2.0, 1.2, 1.8, 1.0])

source_note(slide, "来源：CPCA 2024年全年零售销量数据（品牌口径）")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Competitive Tiers
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "竞争格局分层：\"一超多强\"，新势力加速分化")
page_num(slide, 5)

tiers = [
    ("Tier 1", "绝对龙头", NAVY,
     "BYD\n\n• 市占率 34.1%，销量 427万辆\n• 全价格带覆盖（8万-80万）\n• 垂直整合：电池+芯片+整车",
     "●●●"),
    ("Tier 2", "传统大厂 & 外资", ACCENT,
     "Geely · Tesla · SGMW · Changan\n\n• 合计市占率 ~25%\n• Geely多品牌矩阵\n• Tesla品牌溢价+FSD护城河",
     "●●○"),
    ("Tier 3 (盈利)", "新势力头部", GREEN,
     "Li Auto · AITO\n\n• 理想：增程+家庭定位，毛利20%+\n• 问界：华为智选，智能化领先",
     "●●○"),
    ("Tier 4", "新势力挑战者", ORANGE,
     "Leapmotor · Xpeng · NIO\n\n• 蔚来：换电生态+高端品牌\n• 小鹏：自研XNGP智驾\n• 零跑：性价比+出海",
     "●○○"),
]
for i, (tier, sub, color, body, rating) in enumerate(tiers):
    lx = 0.3 + i * 3.25
    # Tier header box
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.2), Inches(3.1), Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(slide, lx+0.1, 1.22, 2.9, 0.35,
                f"{tier}  {sub}", size=13, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    # Rating
    add_textbox(slide, lx+0.1, 1.6, 2.9, 0.35, f"护城河: {rating}",
                size=12, color=color, align=PP_ALIGN.CENTER)
    # Body card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(2.05), Inches(3.1), Inches(4.85))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xCC,0xD8,0xEE)
    add_textbox(slide, lx+0.12, 2.12, 2.86, 4.7, body,
                size=12, color=DARK_GRAY, wrap=True)

source_note(slide, "分层依据：2024年市占率、盈利状况及护城河深度综合判断")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Technology Differentiation
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "差异化策略：技术路线分化，智能化决胜下半场")
page_num(slide, 6)

data = [
    ["品牌",        "动力路线",     "智能驾驶",      "目标市场",     "价格带（万元）"],
    ["BYD",        "纯电+插混",    "DiPilot（华为合作）","大众+新能源", "8–80"],
    ["Tesla",      "纯电",         "FSD端到端视觉",  "中高端",       "25–55"],
    ["Li Auto",    "增程+纯电",    "AD Max（地平线）","家庭SUV",      "20–50"],
    ["AITO/华为",  "增程+纯电",    "ADS 2.0 领先",   "中高端",       "25–60"],
    ["Xpeng",      "纯电",         "XNGP自研端到端", "中端+科技",    "15–35"],
    ["NIO",        "纯电+换电",    "NIO Pilot",      "高端",         "30–90+"],
    ["Leapmotor",  "纯电",         "基础辅助驾驶",   "性价比",       "8–20"],
    ["Geely/极氪", "纯电+混动",    "Mobileye/自研",  "大众+高端",    "10–60"],
]
add_table(slide, 0.4, 1.2, 12.5, 5.9, data,
          col_widths=[1.8, 1.8, 2.4, 2.0, 2.2])
source_note(slide, "来源：各公司官网、投资者演示材料（2024-2025年）")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Business Model
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "商业模式：四大流派并存，垂直整合与生态赋能领跑")
page_num(slide, 7)

models = [
    ("垂直整合派", NAVY, "BYD",
     "• 自研电池(刀片)、电机、电控、芯片\n"
     "• 成本优势较竞品低 15-20%\n"
     "• 利润护城河最深，规模效应明显\n"
     "• 风险：智能化相对滞后"),
    ("生态赋能派", ACCENT, "AITO（华为智选）",
     "• 华为提供鸿蒙OS+ADS智驾+渠道\n"
     "• 赛力斯专注硬件制造\n"
     "• 品牌高度依赖华为生态\n"
     "• 优势：快速获取华为4000+门店"),
    ("场景专注派", GREEN, "Li Auto（理想）",
     "• 增程式聚焦家庭大型SUV\n"
     "• 毛利率 20%+，新势力唯一盈利\n"
     "• 管理层执行力强，交付高效\n"
     "• 风险：纯电转型压力大"),
    ("服务订阅派", ORANGE, "NIO（蔚来）",
     "• 换电网络+BaaS(电池即服务)模式\n"
     "• NIO House用户社区运营\n"
     "• 高净值用户粘性强\n"
     "• 风险：换电基建成本高，亏损持续"),
]
for i, (title, color, brand, body) in enumerate(models):
    lx = 0.3 + i * 3.25
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.2), Inches(3.1), Inches(0.55))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(slide, lx+0.08, 1.22, 2.94, 0.4,
                f"{title} · {brand}", size=12, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.8), Inches(3.1), Inches(5.1))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xCC,0xD8,0xEE)
    add_textbox(slide, lx+0.12, 1.88, 2.86, 4.9, body,
                size=12, color=DARK_GRAY, wrap=True)

source_note(slide, "商业模式分类基于公司战略定位及收入结构分析")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Moat Framework
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "护城河评估框架：五维度系统化评级")
page_num(slide, 8)

# Left: framework explanation
dims = [
    ("📦 规模经济", "产能利用率、采购议价权、单车制造成本"),
    ("🔬 技术壁垒", "电池技术、自动驾驶算法、芯片自研能力"),
    ("🌟 品牌溢价", "用户NPS、溢价空间、高端化可能性"),
    ("🏪 渠道网络", "门店覆盖度、服务体验、线上销售能力"),
    ("🔗 生态系统", "软件/OTA收入、第三方应用、换电/能源网络"),
]
add_textbox(slide, 0.4, 1.3, 4.5, 0.4,
            "评估维度与考察指标", size=14, bold=True, color=NAVY)
for i, (dim, desc) in enumerate(dims):
    ty = 1.8 + i * 1.05
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(ty), Inches(4.6), Inches(0.95))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = (NAVY if i==0 else
        [ACCENT, GREEN, ORANGE, RED][i-1])
    row_bg.line.fill.background()
    add_textbox(slide, 0.52, ty+0.05, 4.4, 0.4,
                dim, size=13, bold=True, color=WHITE)
    add_textbox(slide, 0.52, ty+0.48, 4.4, 0.4,
                desc, size=11, color=WHITE)

# Right: rating legend
add_textbox(slide, 5.5, 1.3, 7.5, 0.4,
            "评分标准", size=14, bold=True, color=NAVY)
legend_data = [
    ["评级", "说明", "典型特征"],
    ["●●●  强", "难以复制的结构性优势", "行业领先 / 多年积累 / 高切换成本"],
    ["●●○  中", "有一定壁垒但可被追赶", "中等差异化 / 需持续投入维持"],
    ["●○○  弱", "短期竞争优势，护城河浅", "易被复制 / 需监控竞争动态"],
]
add_table(slide, 5.5, 1.8, 7.5, 2.2, legend_data,
          col_widths=[1.5, 2.5, 3.5])

add_textbox(slide, 5.5, 4.2, 7.5, 0.4,
            "护城河综合评级将综合五个维度加权打分，权重：\n"
            "规模经济 25%  |  技术壁垒 30%  |  品牌溢价 20%\n"
            "渠道网络 15%  |  生态系统 10%",
            size=12, color=DARK_GRAY, wrap=True)

source_note(slide, "框架参考：Michael Porter 竞争优势理论 + Morningstar 护城河评级体系")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Moat Ratings (bar chart)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "护城河评级：BYD垂直整合最强，新势力智能化奋起直追")
page_num(slide, 9)

# Clustered bar chart: 5 brands x 5 dimensions (score 1-3)
cd = CategoryChartData()
cd.categories = ["BYD", "Tesla", "Li Auto", "AITO", "NIO", "Xpeng"]
scores = {
    "规模经济":  [3, 2, 2, 1, 1, 1],
    "技术壁垒":  [3, 3, 2, 3, 2, 3],
    "品牌溢价":  [2, 3, 2, 2, 3, 2],
    "渠道网络":  [3, 2, 2, 3, 2, 2],
    "生态系统":  [2, 2, 2, 3, 3, 2],
}
series_colors = [NAVY, ACCENT, GREEN, ORANGE, RED]
for dim, vals in scores.items():
    cd.add_series(dim, vals)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.3), Inches(1.2), Inches(8.5), Inches(5.9), cd
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = True
chart.has_title = False
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.font.size = Pt(9)
for i, rgb in enumerate(series_colors):
    set_series_color(chart, i, rgb)

# Summary table
sum_data = [
    ["品牌", "规模", "技术", "品牌", "渠道", "生态", "综合"],
    ["BYD",   "●●●","●●●","●●○","●●●","●●○","●●●"],
    ["Tesla", "●●○","●●●","●●●","●●○","●●○","●●●"],
    ["Li Auto","●●○","●●○","●●○","●●○","●●○","●●○"],
    ["AITO",  "●○○","●●●","●●○","●●●","●●●","●●○"],
    ["NIO",   "●○○","●●○","●●●","●●○","●●●","●●○"],
    ["Xpeng", "●○○","●●●","●●○","●●○","●●○","●●○"],
]
add_table(slide, 9.0, 1.2, 4.1, 5.9, sum_data,
          col_widths=[1.2, 0.6, 0.6, 0.6, 0.6, 0.6, 0.9])

source_note(slide, "护城河评级为定性分析，●●● =强/3分  ●●○ =中/2分  ●○○ =弱/1分")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – BYD vs Tesla deep dive
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "护城河详解：BYD vs Tesla — 两种路径，各有所长")
page_num(slide, 10)

# BYD column
add_textbox(slide, 0.4, 1.25, 6.0, 0.45,
            "🔵 BYD — 垂直整合之王", size=15, bold=True, color=NAVY)
byd_data = [
    ["指标",          "数值"],
    ["2024年销量",    "427万辆（+41% YoY）"],
    ["市占率",        "34.1%"],
    ["毛利率",        "~18%（汽车业务）"],
    ["净利润",        "~300亿元（预估）"],
    ["核心电池技术",  "刀片电池（磷酸铁锂）"],
    ["自研芯片",      "e平台3.0 + 海鸥芯片"],
    ["价格覆盖",      "8万 – 80万元"],
]
add_table(slide, 0.4, 1.75, 6.0, 3.5, byd_data, col_widths=[2.5, 3.5])

add_textbox(slide, 0.4, 5.35, 6.0, 1.8,
            "✅ 核心优势：刀片电池安全+低成本；全产业链自控；渠道下沉县级市\n"
            "⚠️ 主要风险：品牌向上遭遇天花板；智能驾驶仍落后华为/特斯拉；海外扩张受壁垒",
            size=11.5, color=DARK_GRAY, wrap=True)

# Tesla column
add_textbox(slide, 6.9, 1.25, 6.0, 0.45,
            "🔴 Tesla — 技术品牌双护城河", size=15, bold=True, color=RED)
tesla_data = [
    ["指标",          "数值"],
    ["2024年中国销量","66万辆（-3% YoY）"],
    ["全球市占率",    "~18%（纯电）"],
    ["毛利率",        "~17%（汽车业务）"],
    ["FSD收入潜力",   "每辆车 $8,000-15,000"],
    ["核心技术",      "FSD端到端 + Dojo超算"],
    ["超级充电桩",    "全球7万+站，中国1,500+"],
    ["价格覆盖",      "25万 – 55万元"],
]
add_table(slide, 6.9, 1.75, 6.0, 3.5, tesla_data, col_widths=[2.5, 3.5])

add_textbox(slide, 6.9, 5.35, 6.0, 1.8,
            "✅ 核心优势：FSD软件收入潜力巨大；品牌溢价全球最强；超充网络生态锁定\n"
            "⚠️ 主要风险：中国市场竞争加剧；国内产品迭代速度落后；价格战损毛利",
            size=11.5, color=DARK_GRAY, wrap=True)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.6), Inches(1.2), Inches(0.04), Inches(6.1))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div.line.fill.background()

source_note(slide, "来源：Tesla 2024 Q4 Earnings、BYD 2024年报及公开数据")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – New Forces deep dive
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "护城河详解：新势力分化加剧，理想盈利领跑")
page_num(slide, 11)

nf = [
    ("理想汽车 Li Auto", GREEN,
     "• 2024年交付50万辆（+33%），首个盈利新势力\n"
     "• 增程式家庭SUV定位清晰，避免纯电焦虑\n"
     "• 毛利率 20%+，自由现金流转正\n"
     "• 挑战：L6以下市场、纯电MEGA受挫",
     "●●○", "盈利能力 & 执行力"),
    ("问界 AITO（华为）", ACCENT,
     "• 2024年交付41万辆（+182%），增速第一\n"
     "• 鸿蒙OS + ADS 2.0城区NOA领先行业\n"
     "• 华为渠道赋能4,000+门店覆盖\n"
     "• 挑战：高度依赖华为，品牌独立性弱",
     "●●○", "智能化 & 渠道"),
    ("蔚来 NIO", ORANGE,
     "• 2024年交付31万辆（+39%）\n"
     "• 换电网络2,400站+，BaaS电池订阅模式\n"
     "• 高端圈层用户社区黏性强（NPS领先）\n"
     "• 挑战：持续亏损，换电基建成本高",
     "●●○", "服务生态 & 品牌"),
    ("小鹏汽车 Xpeng", RED,
     "• 2024年交付34万辆（+34%），MONA M03爆款\n"
     "• 自研XNGP端到端智驾，技术实力强\n"
     "• 大众汽车战略入股（4.99%），获技术授权收入\n"
     "• 挑战：品牌溢价弱，中端市场竞争激烈",
     "●●○", "智能驾驶技术"),
]
for i, (name, color, body, rating, strength) in enumerate(nf):
    lx = 0.3 + i * 3.25
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.2), Inches(3.1), Inches(0.55))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(slide, lx+0.08, 1.22, 2.94, 0.4,
                name, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, lx+0.08, 1.82, 2.94, 0.35,
                f"护城河: {rating}  |  优势: {strength}",
                size=10, color=color, align=PP_ALIGN.CENTER)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(2.25), Inches(3.1), Inches(4.65))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xCC,0xD8,0xEE)
    add_textbox(slide, lx+0.12, 2.32, 2.86, 4.5, body,
                size=12, color=DARK_GRAY, wrap=True)

source_note(slide, "来源：各公司月度/季度交付数据及官方投资者文件（2024年）")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Future Trends
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "2025-2027年核心趋势：智能化、出海、整合三大主线")
page_num(slide, 12)

trends = [
    ("趋势一", "价格战趋缓，价值战开启", NAVY,
     "2024-2025年主要价格战触底\n→ 2025H2起竞争转向智能化服务\n→ NOA城区/高速覆盖成营销核心\n→ OTA软件升级成差异化手段"),
    ("趋势二", "智能化决胜：端到端 NOA", ACCENT,
     "华为ADS / 小鹏XNGP / Tesla FSD\n→ 2025年城区NOA覆盖全面铺开\n→ 算法+数据+算力三要素缺一不可\n→ 无自研能力者面临差异化危机"),
    ("趋势三", "出海加速：东南亚+欧洲", GREEN,
     "BYD泰国工厂/欧洲匈牙利工厂\n→ 零跑与Stellantis欧洲合资\n→ 东南亚 2024年中国品牌份额 40%+\n→ 欧盟关税壁垒倒逼本地化生产"),
    ("趋势四", "整合加速，弱者出局", ORANGE,
     "预计2025-2026年有2-3家新势力\n资金耗尽或被并购\n→ CR5市占率预计从58%升至65%+\n→ 国资背景车企并购是主要路径"),
    ("趋势五", "电池技术迭代", RED,
     "固态电池：2026-2027量产预期\n→ 能量密度提升30-50%\n→ BYD/宁德时代/比亚迪核心卡位\n→ 电池成本持续下降利好渗透率"),
    ("趋势六", "商用车+机器人协同", RGBColor(0x6A,0x1B,0x9A),
     "Tesla Optimus机器人 / BYD商用车\n→ 汽车+机器人制造协同效应\n→ 2027年人形机器人开始规模量产\n→ 整车厂具备制造基础优势"),
]
for i, (num, title, color, body) in enumerate(trends):
    row, col = divmod(i, 3)
    lx = 0.3 + col * 4.35; ty = 1.25 + row * 3.05
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(ty), Inches(4.1), Inches(0.6))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(slide, lx+0.1, ty+0.05, 3.9, 0.5,
                f"{num}：{title}", size=12.5, bold=True,
                color=WHITE)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(ty+0.62), Inches(4.1), Inches(2.35))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xCC,0xD8,0xEE)
    add_textbox(slide, lx+0.12, ty+0.7, 3.86, 2.2, body,
                size=11.5, color=DARK_GRAY, wrap=True)

source_note(slide, "趋势预判基于公开研究报告、专家访谈及行业数据综合分析，不构成投资建议")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Market Share Forecast (bar chart)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "2027年市占率预测：格局深化，BYD稳中有降，新势力洗牌")
page_num(slide, 13)

brands_f = ["BYD", "Geely", "Tesla", "SGMW", "Changan", "Li Auto", "AITO", "Others"]
share_24 = [34.1, 7.9, 6.0, 5.9, 5.2, 4.6, 3.8, 32.5]
share_27 = [30.0, 9.0, 5.5, 4.5, 6.0, 7.0, 6.0, 32.0]

cd = CategoryChartData()
cd.categories = brands_f
cd.add_series("2024年实际市占率(%)", share_24)
cd.add_series("2027年预测市占率(%)", share_27)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.3), Inches(1.2), Inches(8.5), Inches(5.9), cd
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = True
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.number_format = '0.0"%"'
chart.plots[0].data_labels.font.size = Pt(9)
set_series_color(chart, 0, NAVY)
set_series_color(chart, 1, ACCENT)

# Commentary
commentary = [
    ("BYD", "稳中微降", "规模仍领先，但竞争加剧和出海摊薄国内份额"),
    ("Geely", "稳步上升", "多品牌矩阵(极氪/领克)持续放量，Flyme车机赋能"),
    ("Li Auto", "显著上升", "纯电L系列+MEGA换代，增量明确"),
    ("AITO", "显著上升", "华为持续赋能，2025款旗舰车型周期开启"),
    ("Tesla", "小幅下降", "中国市场竞争加剧，FSD定价争议"),
    ("SGMW", "下降", "宏光MINI换代迟缓，低价市场被零跑侵蚀"),
]
comm_data = [["品牌", "预测趋势", "核心逻辑"]] + [list(r) for r in commentary]
add_table(slide, 9.0, 1.2, 4.1, 5.9, comm_data, col_widths=[1.0, 1.1, 2.0])

source_note(slide, "预测基于公司指引、产品规划及行业分析综合判断；实际结果可能存在重大偏差")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Conclusion & Investment Implications
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "竞争格局总结：智能化与出海是未来2年最大增量")
page_num(slide, 14)

# 3 column summary
cols = [
    ("🏆 格局结论", NAVY, [
        "BYD：短期无忧，长期需突破智能化瓶颈",
        "Tesla：技术护城河仍深，中国份额承压",
        "理想：新势力最佳盈利模型，转型关键年",
        "问界：华为加持短期红利大，独立性待观察",
        "NIO/Xpeng：技术有亮点，持续盈利待证明",
        "行业整体：2025-2026年将出现2-3家出局",
    ]),
    ("💡 投资启示", ACCENT, [
        "短期：关注盈利可持续性与现金流健康度",
        "中期：重点追踪NOA渗透率与OTA收入占比",
        "长期：出海能力（本地化制造+品牌）是溢价来源",
        "风险点：固态电池量产时点可能重塑竞争格局",
        "优选标的：具备技术+盈利+出海三重优势者",
        "规避：资金储备不足且无差异化优势的尾部品牌",
    ]),
    ("⚡ 关键催化剂", GREEN, [
        "智能驾驶：2025年城区NOA全面商业化落地",
        "出海：东南亚+欧洲中国车市场份额数据",
        "电池：固态电池量产时间表公告",
        "政策：欧盟/美国关税谈判进展",
        "并购：首家头部新势力收并购交易发生",
        "盈利：蔚来/小鹏首次季度盈利时间点",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    lx = 0.3 + i * 4.35
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.25), Inches(4.1), Inches(0.55))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(slide, lx+0.1, 1.27, 3.9, 0.45,
                title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.84), Inches(4.1), Inches(5.35))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xCC,0xD8,0xEE)
    body = "\n".join(f"• {it}" for it in items)
    add_textbox(slide, lx+0.12, 1.92, 3.86, 5.15, body,
                size=12, color=DARK_GRAY, wrap=True)

source_note(slide, "⚠️ 本分析仅供参考和学习，不构成投资建议。投资决策需自行研判并承担相应风险。")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – Risk Factors
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
title_bar(slide, "风险提示：技术、市场、地缘三大风险不可忽视")
page_num(slide, 15)

risks = [
    ("⚙️ 技术风险", RED, [
        "智能驾驶技术迭代不及预期，导致竞争格局骤变",
        "固态电池量产延迟，打乱各厂产品规划节奏",
        "芯片供应受限（地缘政治风险），成本上升",
        "软件安全/数据隐私事件引发监管介入",
    ]),
    ("📊 市场风险", ORANGE, [
        "宏观经济下行影响消费信心和高端车需求",
        "价格战重新激烈，整体毛利率持续承压",
        "电池原材料（锂/钴）价格波动影响成本",
        "NEV渗透率增速放缓（基数效应+潜在买家收窄）",
    ]),
    ("🌐 地缘政治风险", NAVY, [
        "欧盟关税壁垒升级，出海成本大幅提升",
        "中美贸易摩擦延伸至汽车供应链",
        "东南亚/欧洲本土反弹情绪影响中国品牌扩张",
        "华为车BU业务受美国制裁影响核心芯片供应",
    ]),
]
for i, (title, color, items) in enumerate(risks):
    lx = 0.3 + i * 4.35
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.25), Inches(4.1), Inches(0.55))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(slide, lx+0.1, 1.27, 3.9, 0.45,
                title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(1.84), Inches(4.1), Inches(5.35))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xCC,0xD8,0xEE)
    body = "\n".join(f"• {it}" for it in items)
    add_textbox(slide, lx+0.12, 1.92, 3.86, 5.15, body,
                size=12, color=DARK_GRAY, wrap=True)

add_textbox(slide, 0.3, 7.05, 12.7, 0.35,
            "⚠️ 本报告仅供参考和学习，不构成任何形式的投资建议。分析结论基于公开信息，不代表任何机构立场。",
            size=9, color=MID_GRAY, wrap=False)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out_path = "/home/ubuntu/.openclaw-public/workspace/中国新能源汽车行业竞争格局分析_v2.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
